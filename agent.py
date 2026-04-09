"""
SAP Presales Discovery Agent

Uses Claude Opus with web search to research a company and industry,
then generates targeted discovery questions to position SAP solutions.
"""

import io
import os
from pathlib import Path
from typing import Generator
import anthropic

SOLUTIONS_DIR = Path(__file__).parent / "solutions"

# Model ID — override via ANTHROPIC_MODEL env var.
# Direct Anthropic API default: claude-opus-4-6
# Hyperspace proxy default:     anthropic--claude-4.6-opus
MODEL_ID = os.environ.get("ANTHROPIC_MODEL", "claude-opus-4-6")

# Web search is supported on the direct Anthropic API but not through the
# Hyperspace proxy. Enabled automatically when ANTHROPIC_BASE_URL is not set.
_WEB_SEARCH_TOOL = [{"type": "web_search_20260209"}]
_USE_WEB_SEARCH = not bool(os.environ.get("ANTHROPIC_BASE_URL"))

SOLUTION_FILES = {
    "SAP Cloud ERP (Public Cloud)":     "cloud_erp_public.md",
    "SAP Cloud ERP (Private Cloud)":    "cloud_erp_private.md",
    "SAP Business Network":             "sap_business_network.md",
    "SAP Ariba":                        "ariba.md",
    "SAP Integrated Business Planning": "ibp.md",
    "SAP Digital Manufacturing Cloud":  "dmc.md",
    "SAP Field Service Management":     "fsm.md",
}

SYSTEM_PROMPT_TEMPLATE = """You are an expert SAP presales solution advisor with deep knowledge of SAP's portfolio and enterprise software sales methodology. Your role is to help presales professionals prepare for customer discovery calls by:

1. Researching the target company's business context, recent developments, strategic priorities, financial position, and known challenges
2. Analyzing relevant industry trends, pressures, and transformation drivers
3. Generating powerful, open-ended discovery questions specifically designed to surface pain points where the SAP solutions being positioned can deliver value

## SAP SOLUTION CAPABILITIES

The following solutions are being positioned for this customer. Use this knowledge to craft questions that will naturally lead the conversation toward these capabilities:

{solution_capabilities}

## OUTPUT FORMAT

You must produce a structured discovery brief in Markdown with exactly these sections:

---

# Discovery Brief: {{COMPANY}} | {{SOLUTIONS}}

## Company Intelligence
*Key findings from your research — financial health, recent news, strategic initiatives, known challenges, digital transformation status, leadership changes, etc.*

## Industry Context
*Relevant trends, disruptions, regulatory pressures, competitive dynamics, and transformation drivers specific to this industry right now.*

## Strategic Positioning
*Based on company and industry research, your assessment of where SAP has the strongest opportunity and why. Be specific about which pain points are most likely present.*

## Discovery Questions

Organise questions into 3–5 themed sections, each anchored to a **specific business process** where pain is likely to exist based on the company research. Within each theme, write a lead question plus 1–2 natural follow-on probes — questions you'd ask depending on what the customer says.

Use this format:
### Theme: [Business process, e.g., Supplier Order Collaboration]
**Q: [Lead question]**
*Listen for:* [What signal you're hunting and which SAP capability it opens]

> *If they describe manual or disconnected steps:* [Follow-on probe]
> *If they mention exceptions or volume:* [Follow-on probe]

---

## QUESTION QUALITY GUIDELINES

### Anchor to business processes, not pain categories
Each theme must map to a named business process step — not a vague topic like "supply chain" or "reporting". Think in terms of process handoffs:
- Order-to-Cash: customer order intake → fulfilment → invoicing → cash application
- Procure-to-Pay: requisition → sourcing → PO → goods receipt → invoice → payment
- Supply chain collaboration: forecast sharing → order confirmation → supplier commit → exception handling → delivery
- Manufacturing: demand signal → production planning → shop floor execution → quality → goods issue
- Financial close: sub-ledger posting → intercompany → consolidation → reporting

Pick the process steps most relevant to the solutions being positioned and the company context.

### Make questions "process-curious", not accusatory
You are exploring how the process works today — not assuming it's broken. The customer's description of their own process will reveal the inefficiency, without you having to assert it.

- Don't assert pain: "How much time do you waste reconciling supplier invoices?"
- Be process-curious: "Walk me through how a supplier invoice moves from receipt to payment in your organisation — what does that handoff look like today?"

The customer's answer tells you whether there's pain. If there is, you probe deeper. If not, you move on.

### Calibrate specificity — "casually specific"
Questions should be specific enough to sound credible and expert, but not so narrow that if your assumption is slightly off, the question becomes irrelevant. The goal is to stay in the conversation regardless of the answer.

- Too specific (easy to deflect): "How do you handle EDI-based ASN discrepancies with your Tier 1 suppliers?"
- Too vague (no signal): "How's your supply chain working?"
- Casually specific (stays relevant): "When a supplier confirms an order differently from what you sent — different quantity, different date — how does your team find out, and what happens next?"

### Follow-on probes (the "prying" layer)
Within each theme, include 1–2 conditional probes — questions you'd naturally ask if the lead question surfaces something interesting. These are not meant to be asked sequentially; they're options depending on what you hear.

Format them as:
> *If they describe a manual process:* "How often does that create problems downstream — what's the most common failure point?"
> *If they mention it works but is slow:* "What does the delay cost you in practice — is it a working capital issue, a service issue, or something else?"

### SPIN across the thread
Structure each theme's lead question + probes as a natural SPIN arc:
- **Lead question** = Situation or Problem (understand the process, then surface friction)
- **First probe** = Implication (if they confirm pain, expand the consequence)
- **Second probe** = Need-Payoff (get them to articulate what better looks like)

### Volume and selection
- Write exactly **10–12 lead questions** plus follow-on probes within each theme
- Choose 3–5 process themes based on what the research suggests are the most likely pain areas
- Every theme must be justified by something in the company or industry research — not just "this solution covers it"

### Conversational quality
- Lead questions should sound like something you'd say in the first 20 minutes of a meeting with a VP or Director
- Never mention SAP, any competitor, or any product name
- One idea per question — no compound questions
- Use the customer's operational language, not ERP terminology
"""

PROCESS_FLOW_INSTRUCTIONS = """

## PROCESS FLOW DIAGRAM

At the very end of your output — after all discovery questions — append a fenced code block with the tag `process-flow` that describes the key business processes this discovery covers. Example:

```process-flow
{
  "title": "Order-to-Cash Process Flow",
  "processes": [
    {"label": "Customer Order", "action_to_next": "SUBMIT"},
    {"label": "Order Management", "action_to_next": "VALIDATE"},
    {"label": "Fulfillment & Shipping", "action_to_next": "SHIP"},
    {"label": "Invoicing & AR", "action_to_next": "COLLECT"},
    {"label": "Cash Application", "action_to_next": null}
  ]
}
```

Rules:
- Choose 5–7 stages that represent the **customer's** operational process (not SAP modules)
- Labels: 2–4 words, title-case
- `action_to_next`: one uppercase verb (e.g. PLAN, APPROVE, SHIP, INVOICE, COLLECT) — `null` for the final stage
- The title should describe the end-to-end process being discussed (e.g. "Procure-to-Pay Process", "Integrated Supply Chain Flow")
- Base the process on the solutions being positioned and the customer's industry context
"""

OUTPUT_FORMAT_REMINDER = """
Please research this company and industry thoroughly using web search, then produce the complete discovery brief as specified. Make sure questions are specific to this company's situation based on what you find in your research — not generic questions that could apply to any company. Remember to include the process-flow JSON block at the very end.
"""


def extract_file_text(uploaded_file) -> tuple[str, str]:
    """
    Extract plain text from an uploaded Streamlit file object.
    Returns (filename, extracted_text). Returns empty string on failure.
    """
    name = uploaded_file.name
    raw = uploaded_file.read()

    try:
        if name.lower().endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            pages = [page.extract_text() or "" for page in reader.pages]
            return name, "\n\n".join(pages).strip()

        if name.lower().endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(raw))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return name, "\n\n".join(paragraphs).strip()

        if name.lower().endswith(".xlsx"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
            return name, "\n".join(rows).strip()

        if name.lower().endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return name, "\n\n".join(slides).strip()

        # Plain text variants (txt, md, csv, etc.)
        return name, raw.decode("utf-8", errors="replace").strip()

    except Exception as exc:
        return name, f"(Could not extract text: {exc})"


def load_solution_capabilities(solutions: list[str]) -> str:
    """Load and concatenate solution knowledge files for the selected solutions."""
    if not solutions:
        return "No specific solutions selected — generate broad discovery questions."

    capabilities = []
    for solution in solutions:
        filename = SOLUTION_FILES.get(solution)
        if filename:
            file_path = SOLUTIONS_DIR / filename
            if file_path.exists():
                capabilities.append(file_path.read_text(encoding="utf-8"))
            else:
                capabilities.append(f"# {solution}\n(Knowledge file not found — use general knowledge of this solution.)")

    return "\n\n---\n\n".join(capabilities) if capabilities else "No solution files found."


def run_discovery(
    company_name: str,
    industry: str,
    solutions: list[str],
    additional_context: str = "",
    account_context: str = "",
    lob: list[str] | None = None,
    uploaded_files: list | None = None,
) -> Generator[str, None, None]:
    """
    Run the discovery agent and stream the markdown output.

    Yields text chunks as they arrive from the Claude streaming API.
    Handles pause_turn (server-side tool iteration limit) by continuing the loop.
    """
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        yield "**Error:** `ANTHROPIC_API_KEY` environment variable is not set."
        return

    base_url = os.environ.get("ANTHROPIC_BASE_URL")
    client = anthropic.Anthropic(
        api_key=api_key,
        **({"base_url": base_url} if base_url else {}),
    )

    # Build system prompt with injected solution knowledge
    solution_capabilities = load_solution_capabilities(solutions)
    solutions_label = ", ".join(solutions) if solutions else "General SAP Portfolio"
    system_prompt = SYSTEM_PROMPT_TEMPLATE.format(
        solution_capabilities=solution_capabilities
    ) + PROCESS_FLOW_INSTRUCTIONS

    # Extract text from any uploaded files
    files_block = ""
    if uploaded_files:
        sections = []
        for uf in uploaded_files:
            name, text = extract_file_text(uf)
            if text:
                sections.append(f"### {name}\n{text}")
        if sections:
            files_block = "\n\n## Customer-Provided Documents\n\n" + "\n\n---\n\n".join(sections)

    # Build the user prompt
    context_block = f"\n\n**Additional Context Provided:**\n{additional_context.strip()}" if additional_context.strip() else ""
    account_block = f"\n\n## Deal Intelligence (Subjective)\n{account_context.strip()}" if account_context.strip() else ""
    lob_block = f"\n\n**Line of Business Focus:** {', '.join(lob)}\nConcentrate the process themes and discovery questions on these business functions specifically." if lob else ""

    user_message = f"""Please prepare a discovery brief for the following opportunity:

**Company:** {company_name}
**Industry:** {industry}
**SAP Solutions Being Positioned:** {solutions_label}{context_block}{lob_block}{account_block}{files_block}

Using your knowledge of {company_name}, research and summarise:
- Business model, revenue scale, key segments and geographies
- Known strategic priorities, transformation initiatives, or publicly stated goals
- Recognised challenges or pressures facing the business
- Technology landscape or known ERP/system information
- Competitive position and recent performance indicators

Also draw on your knowledge of current trends, pressures, and transformation drivers in the {industry} industry.

Then generate the complete discovery brief as specified in your instructions.{OUTPUT_FORMAT_REMINDER}"""

    messages = [{"role": "user", "content": user_message}]

    # Agentic loop — handles pause_turn for server-side tool iteration limit
    while True:
        with client.messages.stream(
            model=MODEL_ID,
            max_tokens=8000,
            system=system_prompt,
            messages=messages,
            **({"tools": _WEB_SEARCH_TOOL} if _USE_WEB_SEARCH else {}),
        ) as stream:
            for event in stream:
                # Yield text deltas
                if (
                    event.type == "content_block_delta"
                    and event.delta.type == "text_delta"
                ):
                    yield event.delta.text

            final = stream.get_final_message()

        if final.stop_reason == "pause_turn":
            # Server-side tool hit iteration limit — append and continue
            messages.append({"role": "assistant", "content": final.content})
            messages.append({
                "role": "user",
                "content": "Please continue and complete the discovery brief."
            })
        else:
            # end_turn or other terminal reason
            break


def refine_section(
    section_name: str,
    current_content: str,
    full_brief: str,
    instructions: str,
    company: str,
    industry: str,
    solutions: list[str],
) -> Generator[str, None, None]:
    """
    Stream a rewritten version of a single brief section.

    Yields text chunks. The caller is responsible for splicing the result
    back into the full brief.
    """
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        yield "**Error:** `ANTHROPIC_API_KEY` environment variable is not set."
        return

    base_url = os.environ.get("ANTHROPIC_BASE_URL")
    client = anthropic.Anthropic(
        api_key=api_key,
        **({"base_url": base_url} if base_url else {}),
    )

    solution_capabilities = load_solution_capabilities(solutions)
    system_prompt = SYSTEM_PROMPT_TEMPLATE.format(solution_capabilities=solution_capabilities)

    user_message = f"""You are refining the **{section_name}** section of a discovery brief for {company} ({industry}).

Here is the full discovery brief for context:

<full_brief>
{full_brief}
</full_brief>

Here is the current content of the {section_name} section:

<current_section>
{current_content}
</current_section>

Refinement instructions:
{instructions}

Output ONLY the revised body content for this section. Do NOT include the section heading (e.g., "## {section_name}") — just the body. Preserve all markdown formatting conventions from the original brief (bold questions, italic listen-for lines, blockquote probes, etc.)."""

    with client.messages.stream(
        model="anthropic--claude-4.6-opus",
        max_tokens=4000,
        system=system_prompt,
        messages=[{"role": "user", "content": user_message}],
    ) as stream:
        for event in stream:
            if (
                event.type == "content_block_delta"
                and event.delta.type == "text_delta"
            ):
                yield event.delta.text


# ── Demo Prep ─────────────────────────────────────────────────────────────────

DEMO_PREP_SYSTEM_PROMPT = """You are an expert SAP presales solution advisor preparing a solution consultant or account executive for a customer demo. You have deep knowledge of SAP's product portfolio, enterprise software sales methodology, and what makes a demo compelling vs. forgettable.

Your role is to turn a discovery brief into a concrete demo preparation guide that:
1. Maps the specific pain points uncovered in discovery to the SAP capabilities that address them
2. Builds a realistic, timed demo agenda that flows naturally and stays within the allotted time
3. Writes talking points that reference the customer's own words and situation — not generic SAP marketing language
4. Anticipates the objections most likely to surface based on the account context
5. Identifies the 2–3 moments in the demo where you pause, name the value, and invite the customer to react

## SAP SOLUTION CAPABILITIES

{solution_capabilities}

## OUTPUT FORMAT

Produce a structured demo prep brief in Markdown with exactly these sections:

---

# Demo Prep: {{COMPANY}} | {{SOLUTIONS}}

## Demo Agenda

A timed table mapping each segment to a specific pain point from the discovery brief. Format:

| Time | Segment | Pain Point Addressed |
|------|---------|---------------------|
| 0:00–0:05 | Opening & agenda | Set context, confirm priorities |
| 0:05–0:20 | [Capability area] | [Specific pain from discovery] |
...

Keep segments to 5–8 minutes each. The final 5 minutes should always be "Next Steps & Q&A."

## Talking Points

For each demo segment, write:

### [Segment name]
**What to show:** [The specific screen, flow, or capability — be precise]
**What to say:** [2–3 sentences, written as you'd actually say them. Reference the customer's situation specifically, not generic product descriptions.]
**Connect back to:** [The specific discovery finding this addresses — quote or paraphrase what the customer said or what the research surfaced]

## Anticipated Objections

For each likely objection:

### "[The objection, as the customer would say it]"
**Why they're raising it:** [The underlying concern]
**Response:** [How to acknowledge and redirect — 2–3 sentences]
**If they push further:** [The follow-on response or question to ask]

## Wow Moments

Identify exactly 2–3 moments in the demo where you should pause, name the value, and invite the customer to react. These are not feature highlights — they are moments where the product directly mirrors a specific pain the customer described.

### Wow Moment [N]: [Short name]
**When:** [Which demo segment, approximately which minute]
**Setup:** [1 sentence — prime the customer before showing it: "You mentioned earlier that..."]
**The moment:** [What you show or say]
**Invite reaction:** [The question you ask after: "Does that reflect how it works today?" / "What would that mean for your team?"]

## Value Story

A 3–5 sentence narrative that ties everything together. This is what you'd say if the CFO walked in at the end and asked "so what's the bottom line?" It should connect the specific pain points from discovery to the specific capabilities shown to a concrete business outcome — using the customer's language, not SAP's.

---

## QUALITY GUIDELINES

- Every segment, talking point, and wow moment must be grounded in something specific from the discovery brief — not generic SAP capability descriptions
- Talking points should sound like something you'd say in a conversation, not a product brochure
- Never mention SAP product names, version numbers, or module codes when writing what to say — use business language
- Objection responses must acknowledge the concern genuinely before redirecting
- The value story must be specific enough that it could only apply to this customer
"""

DEMO_DELIVERY_SYSTEM_PROMPT = """You are an expert SAP presales solution advisor helping a solution consultant or account executive prepare for the day of a customer demo. The demo prep has already been done — now you are producing the delivery tools: what to have in hand, what to do immediately after, and how to communicate the outcome to stakeholders.

## SAP SOLUTION CAPABILITIES

{solution_capabilities}

## OUTPUT FORMAT

Produce a demo delivery pack in Markdown with exactly these sections:

---

# Demo Delivery Pack: {{COMPANY}}

## Live Cheat Sheet

A scannable one-pager the presenter can glance at during the demo. Format as short bullets under bold section headers. Include:

**Opening** — 2–3 bullets: the context-setting sentence, the agenda confirmation question, the "permission to probe" framing
**Transitions** — one bullet per demo segment transition: the bridging sentence that connects one segment to the next
**Timing reminders** — key time checkpoints (e.g., "By :30 you should be in Talking Points")
**Things not to forget** — 3–5 bullets: the specific things most likely to be skipped under pressure (the wow moment setup, the question at the end of each segment, the "does this resonate" pause)
**Closing** — the question to ask before going to next steps: something that surfaces a clear signal on where the deal stands

## Post-Demo Follow-Up Plan

**Within 2 hours:**
[Specific actions — what to send, who to copy, what to capture internally]

**Follow-up email draft:**
[A complete, ready-to-send email. Natural tone — not corporate. Subject line included. Body should: thank them for the time (one sentence), summarise the 2–3 things that resonated most (reference specific moments from the demo), state the agreed next step, and close with a door-opener question. 150–200 words maximum.]

**Open questions to close before the next meeting:**
[Bulleted list — the specific things you still need to learn or confirm to advance the deal]

## Executive Summary Slide Outline

An outline for a one-slide or one-page summary for stakeholders who weren't in the room.

**Slide title:** [Short, value-focused — not "SAP Demo Summary"]
**Opening line:** [One sentence that frames the business problem — specific to this customer]
**Three value points:** [Bullet 1: capability shown → outcome it addresses], [Bullet 2], [Bullet 3]
**Call to action:** [What you want the reader to do or approve]
**Supporting evidence:** [Any specific number, benchmark, or customer reference to include]

---

## QUALITY GUIDELINES

- The cheat sheet must be scannable at a glance — no paragraphs, no full sentences except in transitions
- The follow-up email must sound like it was written by a human who had the meeting — reference specific moments, not generic capabilities
- The exec summary must be self-contained — someone who wasn't in the room should be able to understand the value and the next step
- Everything must be grounded in the demo prep brief provided — do not introduce new content
"""


def run_demo_prep(
    company_name: str,
    industry: str,
    solutions: list[str],
    discovery_brief: str,
    demo_duration: str,
    demo_focus: str = "",
    audience: str = "",
) -> Generator[str, None, None]:
    """
    Generate a demo preparation brief from a discovery brief.
    Streams markdown output chunks.
    """
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        yield "**Error:** `ANTHROPIC_API_KEY` environment variable is not set."
        return

    base_url = os.environ.get("ANTHROPIC_BASE_URL")
    client = anthropic.Anthropic(
        api_key=api_key,
        **({"base_url": base_url} if base_url else {}),
    )

    solution_capabilities = load_solution_capabilities(solutions)
    solutions_label = ", ".join(solutions) if solutions else "General SAP Portfolio"
    system_prompt = DEMO_PREP_SYSTEM_PROMPT.format(solution_capabilities=solution_capabilities)

    focus_block = f"\n\n**Demo Focus / Special Instructions:** {demo_focus.strip()}" if demo_focus.strip() else ""
    audience_block = f"\n\n**Audience in the Room:** {audience.strip()}" if audience.strip() else ""

    user_message = f"""Please prepare a demo preparation brief for the following opportunity:

**Company:** {company_name}
**Industry:** {industry}
**SAP Solutions Being Demonstrated:** {solutions_label}
**Demo Duration:** {demo_duration}{focus_block}{audience_block}

## Discovery Brief (Context for this Demo)

{discovery_brief}

Using the discovery brief above, produce the complete demo prep brief as specified in your instructions. Every section must reference specific findings from the discovery brief — not generic SAP content."""

    messages = [{"role": "user", "content": user_message}]

    while True:
        with client.messages.stream(
            model=MODEL_ID,
            max_tokens=8000,
            system=system_prompt,
            messages=messages,
        ) as stream:
            for event in stream:
                if (
                    event.type == "content_block_delta"
                    and event.delta.type == "text_delta"
                ):
                    yield event.delta.text

            final = stream.get_final_message()

        if final.stop_reason == "pause_turn":
            messages.append({"role": "assistant", "content": final.content})
            messages.append({"role": "user", "content": "Please continue and complete the demo prep brief."})
        else:
            break


def run_demo_delivery(
    company_name: str,
    industry: str,
    solutions: list[str],
    demo_prep_brief: str,
    last_minute_context: str = "",
) -> Generator[str, None, None]:
    """
    Generate a demo delivery pack (cheat sheet, follow-up plan, exec summary) from a demo prep brief.
    Streams markdown output chunks.
    """
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        yield "**Error:** `ANTHROPIC_API_KEY` environment variable is not set."
        return

    base_url = os.environ.get("ANTHROPIC_BASE_URL")
    client = anthropic.Anthropic(
        api_key=api_key,
        **({"base_url": base_url} if base_url else {}),
    )

    solution_capabilities = load_solution_capabilities(solutions)
    solutions_label = ", ".join(solutions) if solutions else "General SAP Portfolio"
    system_prompt = DEMO_DELIVERY_SYSTEM_PROMPT.format(solution_capabilities=solution_capabilities)

    context_block = f"\n\n**Last-Minute Context:** {last_minute_context.strip()}" if last_minute_context.strip() else ""

    user_message = f"""Please prepare a demo delivery pack for the following opportunity:

**Company:** {company_name}
**Industry:** {industry}
**SAP Solutions:** {solutions_label}{context_block}

## Demo Prep Brief

{demo_prep_brief}

Using the demo prep brief above, produce the complete demo delivery pack as specified in your instructions."""

    messages = [{"role": "user", "content": user_message}]

    while True:
        with client.messages.stream(
            model=MODEL_ID,
            max_tokens=6000,
            system=system_prompt,
            messages=messages,
        ) as stream:
            for event in stream:
                if (
                    event.type == "content_block_delta"
                    and event.delta.type == "text_delta"
                ):
                    yield event.delta.text

            final = stream.get_final_message()

        if final.stop_reason == "pause_turn":
            messages.append({"role": "assistant", "content": final.content})
            messages.append({"role": "user", "content": "Please continue and complete the demo delivery pack."})
        else:
            break
